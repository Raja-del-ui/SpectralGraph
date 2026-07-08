#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SpectralGraph"
    subtitle.text = "A NoSQL-Backed Red Team Asset & Vulnerability Tracker\n\nZuni\n4th Semester - BSc Computer Science\nAbbottabad University of Science & Technology"

    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem: Manual Tracking"
    tf = slide.placeholders[1].text_frame
    tf.text = "Security professionals generate massive amounts of unstructured data."
    tf.add_paragraph().text = "Spreadsheets lack historical context (overwriting data destroys audit trails)."
    tf.add_paragraph().text = "Flat files cannot handle complex, many-to-many network relationships efficiently."
    tf.add_paragraph().text = "No application-layer access control, creating security risks."

    # Slide 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Solution: SpectralGraph"
    tf = slide.placeholders[1].text_frame
    tf.text = "An automated, NoSQL-backed database application."
    tf.add_paragraph().text = "Automatically ingests raw network scan data (Nmap)."
    tf.add_paragraph().text = "Tracks infrastructure changes over time (Time-Series Logging)."
    tf.add_paragraph().text = "Securely maps CVEs to specific targets using relational references."

    # Slide 4: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technology Stack"
    tf = slide.placeholders[1].text_frame
    tf.text = "Database Engine: MongoDB (NoSQL Document Store)"
    tf.add_paragraph().text = "Backend Logic: Python 3"
    tf.add_paragraph().text = "Integration Libraries: PyMongo, python-nmap"
    tf.add_paragraph().text = "Security: hashlib (SHA-256 Authentication)"
    tf.add_paragraph().text = "Environment: Kali Linux"

    # Slide 5: Database Architecture (Schema)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Database Architecture: Hybrid Schema"
    tf = slide.placeholders[1].text_frame
    tf.text = "We utilized both NoSQL paradigms to optimize the database:"
    tf.add_paragraph().text = "Embedding: `targets` collection embeds open ports directly for rapid reads."
    tf.add_paragraph().text = "Referencing: `scan_history` links to targets to maintain an audit trail."
    tf.add_paragraph().text = "Referencing: `vulnerabilities` normalizes data by linking CVEs to target IDs."
    tf.add_paragraph().text = "Standalone: `users` collection manages credentials securely."

    # Slide 6: Automated Data Ingestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Feature: Automated Ingestion"
    tf = slide.placeholders[1].text_frame
    tf.text = "Replaces manual data entry with Python automation."
    tf.add_paragraph().text = "Idempotent Operations: Uses PyMongo's `upsert=True` logic."
    tf.add_paragraph().text = "Prevents duplicate entries when re-scanning existing infrastructure."
    tf.add_paragraph().text = "Maintains a real-time 'Current State' of the network."

    # Slide 7: Time-Series Audit Logging
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Feature: Time-Series Audit"
    tf = slide.placeholders[1].text_frame
    tf.text = "Security requires knowing *when* a change occurred."
    tf.add_paragraph().text = "The Python backend simultaneously writes to the `scan_history` collection."
    tf.add_paragraph().text = "Creates an immutable log of what ports were open at specific timestamps."
    tf.add_paragraph().text = "Demonstrates One-to-Many temporal relationships."

    # Slide 8: Role-Based Access Control (RBAC)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Feature: Security & RBAC"
    tf = slide.placeholders[1].text_frame
    tf.text = "Application-layer security prevents unauthorized database modifications."
    tf.add_paragraph().text = "Passwords are never stored in plain-text (SHA-256 Hashing)."
    tf.add_paragraph().text = "Role-Based Access Control enforced at the script level."
    tf.add_paragraph().text = "Only 'Lead_Operator' accounts can inject new CVEs; 'Analyst' accounts are blocked."

    # Slide 9: Advanced Aggregation Pipelines
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Feature: Advanced Aggregations"
    tf = slide.placeholders[1].text_frame
    tf.text = "Complex analytics performed natively in the NoSQL engine."
    tf.add_paragraph().text = "Stage 1 ($lookup): Joins targets and vulnerabilities (NoSQL LEFT JOIN)."
    tf.add_paragraph().text = "Stage 2 ($unwind): Deconstructs vulnerability arrays for counting."
    tf.add_paragraph().text = "Stage 3 ($group): Dynamically groups CVEs by Operating System and Severity."
    tf.add_paragraph().text = "Stage 4 ($sort): Orders the output for dashboard presentation."

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion & Live Demonstration"
    tf = slide.placeholders[1].text_frame
    tf.text = "SpectralGraph successfully bridges networking and modern database structures."
    tf.add_paragraph().text = "Ready for live demonstration of:"
    tf.add_paragraph().text = "1. Nmap ingestion into MongoDB."
    tf.add_paragraph().text = "2. Security block on the Analyst account."
    tf.add_paragraph().text = "3. Executing the Aggregation Pipeline dashboard."

    prs.save('SpectralGraph_Presentation.pptx')
    print("[+] Presentation successfully saved as 'SpectralGraph_Presentation.pptx'")

if __name__ == '__main__':
    create_presentation()
